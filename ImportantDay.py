# --------------- Dependencies --------------- 
# pip install requests pywin32 python-pptx   |
# --------------------------------------------

import requests
import json
import csv
import os
import time
import pythoncom
import win32com.client
from io import StringIO
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from xml.etree import ElementTree as ET



# ---------------------- Test Data ----------------------
USE_TEST_DATA = False
today = datetime.today()
today_str = today.strftime("%m/%d/%Y")

test_csv_data = f"""/ 
"Employee Id","First Name","Last Name","Date Birthday","Date Hired","Employee Status","In Payroll","Badge Type"
"9001","Alice","BirthdayOnly","03/04/2000","01/01/2000","Active","Yes","Simple #1"
"9002","Bob","AnniversaryOnly","01/01/1970","06/02/2005","Active","Yes","Simple #1"
"9003","Charlie","BothToday","{today_str}","{today_str}","Active","Yes","Simple #1"
"9003","Charlie","BothToday","{today_str}","{today_str}","Active","Yes","Simple #1"
"""

# ---------------------- Functions ----------------------
DEFAULT_FONT_SIZE = Pt(24)
DEFAULT_FONT_NAME = "Calibri"
DEFAULT_FONT_COLOR = RGBColor(0, 0, 0)

def clean_existing_boxes(slide, prefix):
    for shape in list(slide.shapes):
        if shape.name.startswith(prefix):
            slide.shapes._spTree.remove(shape._element)

def add_textbox(slide, name, text, left, top, width, height):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = name
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = DEFAULT_FONT_SIZE
    font.name = DEFAULT_FONT_NAME
    font.color.rgb = DEFAULT_FONT_COLOR
    return textbox

def run_slide_update(prs):
    slide = prs.slides[7]
    clean_existing_boxes(slide, "BirthdayBox")
    clean_existing_boxes(slide, "AnniversaryBox")
    left_x, right_x = Inches(0), Inches(6.67)
    top_y, spacing = Inches(2.01), Inches(0.85)

    for i, person in enumerate(birthdays):
        text = person.get("__birthday_text", f"{person['First Name']} {person['Last Name']}")
        add_textbox(slide, f"BirthdayBox{i+1}", text, left_x, top_y + i * spacing, Inches(6.67), Inches(1.04))

    for i, person in enumerate(anniversaries):
        years = person.get("Years", "?")
        text = person.get("__anniversary_text", f"{person['First Name']} {person['Last Name']} - {years} Year(s)")
        add_textbox(slide, f"AnniversaryBox{i+1}", text, right_x, top_y + i * spacing, Inches(6.67), Inches(1.04))

def get_shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None

def safe_set_text(slide, shape_name, text):
    shape = get_shape_by_name(slide, shape_name)
    if shape and shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()  # Clear existing content
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = text

        if shape_name == "TemperatureBox":
            run.font.name = 'Arial Rounded MT Bold'
            run.font.size = Pt(72)
            run.font.bold = True
        elif shape_name == "HumidityBox":
            run.font.name = 'Arial Rounded MT Bold'
            run.font.size = Pt(36)
            run.font.bold = True
        elif shape_name == "WindBox":
            run.font.name = 'Arial Rounded MT Bold'
            run.font.size = Pt(20)
            run.font.bold = False
        elif shape_name == "WeatherBox":
            run.font.name = 'Arial Rounded MT Bold'
            run.font.size = Pt(24)
            run.font.bold = False
        else:
            run.font.name = 'Arial Rounded MT Bold'
            run.font.size = Pt(20)
            run.font.bold = False

        if shape_name in ["FactBox", "DateBox"]:
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            if shape_name in "DateBox":
                run.font.size = Pt(36)
        else:
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        run.font.italic = False

        return True
    return False



def update_weather_slide(prs):
    slide = prs.slides[8]

    # Load weather XML
    station_id = "KDAY"
    url = f"https://forecast.weather.gov/xml/current_obs/{station_id}.xml"
    response = requests.get(url)
    if response.status_code != 200:
        print("Failed to fetch weather data.")
        return

    xml_root = ET.fromstring(response.content)

    def get_xml_value(tag_name):
        element = xml_root.find(tag_name)
        return element.text if element is not None else "N/A"

    location = get_xml_value("location")
    temperature = get_xml_value("temperature_string")
    weather = get_xml_value("weather")
    humidity = get_xml_value("relative_humidity") + "%"
    wind = get_xml_value("wind_string")
    obs_time = get_xml_value("observation_time")

    # Hide all groups initially
    for i in range(1, 8):
        group_name = f"Day{i}Group"
        shape = get_shape_by_name(slide, group_name)
        if shape:
            shape.visible = False

    safe_set_text(slide, "WeatherForBox", location)
    safe_set_text(slide, "TimeBox", obs_time)
    safe_set_text(slide, "WeatherBox", weather)
    safe_set_text(slide, "HumidityBox", humidity)
    safe_set_text(slide, "WindBox", wind)
    safe_set_text(slide, "TemperatureBox", temperature)

    update_current_condition_image(slide, weather)

def update_current_condition_image(slide, weather):
    emoji_url = ""
    wx_lower = weather.lower()
    is_daytime = 6 <= datetime.now().hour < 18

    if "hail" in wx_lower or "sleet" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/hail.png"
    elif "snow" in wx_lower or "flurr" in wx_lower or "blizzard" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/snow.png"
    elif "storm" in wx_lower or "thunder" in wx_lower or "lightning" in wx_lower or "tornado" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/thunderstorm.png"
    elif "overcast" in wx_lower or "cast" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/cloudy.png"
    elif "clear" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/sun.png" if is_daytime else "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/moon.png"
    elif "rain" in wx_lower or "showers" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/rain.png"
    elif "cloud" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/dayCloudy.png" if is_daytime else "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/nightCloudy.png"
    elif "fog" in wx_lower or "haze" in wx_lower or "mist" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/dayFog.png" if is_daytime else "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/nightFog.png"
    elif "sun" in wx_lower:
        emoji_url = "C:/Users/User/Documents/vscode/EmployeeImportantDay/media/sun.png"

    # Remove old image
    shape = get_shape_by_name(slide, "CurrentWeatherImage")
    if shape:
        slide.shapes._spTree.remove(shape._element)

    # Add new image
    if emoji_url:
        img = slide.shapes.add_picture(emoji_url, Inches(5.70), Inches(3.29), Inches(1.93), Inches(2.00))
        img.name = "CurrentWeatherImage"

def wait_until_file_unlocked(filepath, timeout=15):
    start_time = time.time()
    while True:
        try:
            with open(filepath, 'a'):
                return True
        except PermissionError:
            if time.time() - start_time > timeout:
                print("Timeout: PowerPoint did not release the file.")
                return False
            time.sleep(0.5)


def update_fact_slide(prs):
    slide = prs.slides[1]
    today = datetime.today()
    month = today.month
    day = today.day

    # Wikipedia API
    url = f"https://en.wikipedia.org/api/rest_v1/feed/onthisday/events/{month:02d}/{day:02d}"
    try:
        response = requests.get(url)
        response.raise_for_status()
        data = response.json()
    except Exception as e:
        print(f"Failed to fetch historical data: {e}")
        return

    if not data.get("events"):
        print("No historical events found.")
        return

    # Pick the most relevant/first event
    event = data["events"][0]
    year = event["year"]
    description = event["text"]

    # Try to get image URL
    image_url = None
    try:
        if event.get("pages"):
            for page in event["pages"]:
                if "originalimage" in page:
                    image_url = page["originalimage"]["source"]
                    break
    except Exception:
        pass

    # Set text in placeholders
    safe_set_text(slide, "DateBox", f"{today.strftime('%B')} {day}, {year}")
    safe_set_text(slide, "FactBox", description)

    # Remove old image
    old_img = get_shape_by_name(slide, "FactImage")
    if old_img:
        slide.shapes._spTree.remove(old_img._element)

    # Add new image from URL
    if image_url:
        try:
            img_response = requests.get(image_url, stream=True)
            if img_response.status_code == 200:
                with open("temp_fact_image.jpg", "wb") as f:
                    f.write(img_response.content)
                img = slide.shapes.add_picture("temp_fact_image.jpg", Inches(9.42), Inches(4.53), Inches(3), Inches(2))
                img.name = "FactImage"
        except Exception as e:
            print(f"Failed to insert image: {e}")





# -------------------- Fetch Data --------------------
if USE_TEST_DATA:
    data = test_csv_data
else:
    loginId = {'credentials': {'username': '', 'company': '', 'password': ''}}
    loginUrl = 'http://secure3.saashr.com/ta/rest/v1/login'
    loginHeaders = {'Content-type': 'application/json', 'api-key': ''}
    login = requests.post(loginUrl, headers=loginHeaders, json=loginId)
    if login.status_code == requests.codes.ok:
        token = json.loads(login.text)["token"]
        sessionHeaders = {'Content-type': 'application/json', 'Authentication': 'bearer ' + token}
        report = requests.get('http://secure3.saashr.com/ta/rest/v1/report/saved/107538439', headers=sessionHeaders)
        data = report.text
    else:
        print("Login failed:", login.text)
        exit()

# -------------------- Parse Data --------------------
csv_reader = csv.DictReader(StringIO(data))
birthdays, anniversaries = [], []

# Determine the start (Monday) and end (Sunday) of this week
today = datetime.today()
start_of_week = today - timedelta(days=today.weekday())  # Monday
end_of_week = start_of_week + timedelta(days=6)          # Sunday

for row in csv_reader:
    try:
        bday = datetime.strptime(row["Date Birthday"], "%m/%d/%Y")
        hire = datetime.strptime(row["Date Hired"], "%m/%d/%Y")
    except ValueError:
        continue

    # Adjust to current year for birthday comparison
    bday_this_year = bday.replace(year=today.year)
    hire_this_year = hire.replace(year=today.year)

    # --- Birthday Logic ---
    if start_of_week.date() <= bday_this_year.date() <= end_of_week.date():
        if bday_this_year.date() == today.date():
            row["__birthday_text"] = f"Happy birthday to {row['First Name']} {row['Last Name']}!"
        else:
            row["__birthday_text"] = f"{row['First Name']} {row['Last Name']}'s birthday is this week!"
        birthdays.append(row)

    # --- Work Anniversary Logic ---
    if start_of_week.date() <= hire_this_year.date() <= end_of_week.date():
        years = today.year - hire.year
        year_label = "year" if years == 1 else "years"
        if hire_this_year.date() == today.date():
            row["__anniversary_text"] = f"{years} {year_label} today! Congrats, {row['First Name']} {row['Last Name']}!"
        else:
            row["__anniversary_text"] = f"{row['First Name']} {row['Last Name']} has been with us for {years} {year_label} this week!"
        row["Years"] = years
        anniversaries.append(row)

# ---------------- Modify PowerPoint ----------------
# A lot of this is likely filler and/or unnecessary, but it is the first thing in two days that has worked so we're gonna leave that there
ppt_path = "C:/Users/User/Documents/vscode/EmployeeImportantDay/DisplayFinal.pptx"
try:
    pythoncom.CoInitialize()
    try:
        ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
        ppt_app.Quit()
        print("Closed PowerPoint via COM.")
        time.sleep(2)
    except Exception:
        print("PowerPoint COM object not found or already closed.")

    os.system("taskkill /f /im POWERPNT.EXE >nul 2>&1")
    os.system("taskkill /f /im OfficeClickToRun.exe >nul 2>&1")
    os.system("taskkill /f /im OfficeBackgroundTaskHandler.exe >nul 2>&1")
    print("Closed running PowerPoint instance.")

    time.sleep(5)

    lockfile = "C:/Users/User/Documents/vscode/EmployeeImportantDay/~$DisplayFinal.pptx"
    if os.path.exists(lockfile):
        os.remove(lockfile)
        print("Deleted leftover lock file.")

    # Extra cleanup
    for filename in os.listdir(os.path.dirname(ppt_path)):
        if filename.startswith("~$") or filename.endswith(".tmp"):
            try:
                os.remove(os.path.join(os.path.dirname(ppt_path), filename))
                print(f"Deleted temp file: {filename}")
            except Exception:
                pass

    time.sleep(5)
    if not wait_until_file_unlocked(ppt_path):
        print("File is still locked. Aborting")
        exit()

except Exception:
    print("PowerPoint was not open or already closed.")

prs = Presentation(ppt_path)
run_slide_update(prs)
update_weather_slide(prs)
update_fact_slide(prs)
# Wait until the target file is no longer locked
while True:
    try:
        with open(ppt_path, "a"):
            break
    except PermissionError:
        print("Waiting for file to be unlocked...")
        time.sleep(1)
prs.save(ppt_path)
print("PowerPoint slides updated.")
# Remove temp image if it exists
if os.path.exists("temp_fact_image.jpg"):
    try:
        os.remove("temp_fact_image.jpg")
        print("Temporary fact image deleted.")
    except Exception as e:
        print(f"Failed to delete temp image: {e}")


# Optional: Reopen PowerPoint to view results
try:
    powerpoint_path = r"C:/Program Files/Microsoft Office/root/Office16/POWERPNT.EXE"
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        # Open the presentation
        presentation = ppt_app.Presentations.Open(ppt_path, WithWindow=True)
        print("PowerPoint opened.")

        # Start the slideshow
        presentation.SlideShowSettings.Run()
        print("Slideshow started.")
    except Exception as e:
        print(f"Failed to open PowerPoint: {e}")
except Exception as e:
    print(f"Failed to reopen PowerPoint: {e}")
